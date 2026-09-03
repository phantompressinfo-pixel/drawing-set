"""Build the Office Hub page deck from the treatment-33 renders."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "office-hub-pages.pptx")

W_IN, H_IN = 13.333, 7.5
NAVY = RGBColor(0x02, 0x20, 0x49)
NAVY_DEEP = RGBColor(0x01, 0x14, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xC6, 0xD6, 0xE8)
GOTHAM = "Gotham"          # falls back gracefully; Google Slides will substitute

PAGES = [
    ("page_01_office.jpg",        "Home",              "The landing page. Search, five one-click actions, and the six sections the office uses most."),
    ("page_02_announcements.jpg", "Announcements",     "Office reminders and notices. Newest first, each one dated and tagged."),
    ("page_03_office.jpg",        "Office Standards",  "Drafting conventions, file naming, sheet setup — the Office Standards Manual and everything under it."),
    ("page_04_templates.jpg",     "Templates",         "Start-here files: sheet sets, letterhead, proposals, presentation decks."),
    ("page_05_forms.jpg",         "Forms",             "Time off, expenses, IT requests — every form the office fills out, in one list."),
    ("page_06_standard.jpg",      "SOPs",              "Step-by-step procedures for how the office runs a project from kickoff to closeout."),
    ("page_07_revit.jpg",         "Revit Standards",   "Model setup, worksharing, families, view templates, annotation."),
    ("page_08_learning.jpg",      "Learning Sessions", "The weekly presentation program. Upcoming sessions, the sign-up, and the full archive by month."),
]


def txt(slide, x, y, w, h, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        spacing=0, font=GOTHAM):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    if spacing:
        r.font._rPr.set("spc", str(int(spacing * 100)))
    return tb


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


prs = Presentation()
prs.slide_width = Inches(W_IN)
prs.slide_height = Inches(H_IN)
blank = prs.slide_layouts[6]

# ---- title slide -------------------------------------------------
s = prs.slides.add_slide(blank)
bg(s, NAVY)
s.shapes.add_shape(1, Inches(1.05), Inches(2.62), Inches(0.055), Inches(1.55)).fill.solid()
bar = s.shapes[-1]
bar.fill.fore_color.rgb = RGBColor(0x6C, 0x94, 0xC4)
bar.line.fill.background()
bar.shadow.inherit = False
txt(s, 1.42, 2.58, 10.5, 0.5, "EIGELBERGER ARCHITECTURE & DESIGN", 13, True, ICE, spacing=2.6)
txt(s, 1.42, 3.05, 10.5, 0.9, "Office Hub", 46, True, WHITE)
txt(s, 1.45, 4.02, 9.6, 0.6, "Every page, in the proposed design.", 17, False, ICE)
txt(s, 1.05, 6.55, 8.0, 0.4, "Built in Google Sites  ·  files live in Drive  ·  no page edits to keep it current",
    12, False, RGBColor(0x8F, 0xAA, 0xC9))

# ---- one slide per page ------------------------------------------
for fn, title, note in PAGES:
    path = os.path.join(HERE, fn)
    assert os.path.exists(path), path
    s = prs.slides.add_slide(blank)
    bg(s, NAVY_DEEP)
    s.shapes.add_picture(path, 0, 0, Inches(W_IN), Inches(H_IN))
    # speaker note carries the description so the image stays full-bleed
    s.notes_slide.notes_text_frame.text = "%s — %s" % (title, note)

prs.save(OUT)
print("saved", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
